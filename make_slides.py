"""Generate SaaSBench.pptx — a Google-Slides-importable pitch deck (8 slides).
Run:  python3 make_slides.py   ->   SaaSBench.pptx
Upload to Drive and open with Google Slides (or Slides → File → Import slides)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette (matches the web deck) ----
BG     = RGBColor(0x0A, 0x0B, 0x0F)
CARD   = RGBColor(0x15, 0x18, 0x20)
TRACK  = RGBColor(0x0E, 0x10, 0x18)
BORDER = RGBColor(0x2A, 0x2E, 0x40)
TEXT   = RGBColor(0xD6, 0xD8, 0xDE)
SOFT   = RGBColor(0x9A, 0xA0, 0xB0)
DIM    = RGBColor(0x6B, 0x70, 0x84)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACC    = RGBColor(0x7C, 0x8B, 0xFF)   # purple
GRN    = RGBColor(0x38, 0xEF, 0x7D)   # green
BLUE   = RGBColor(0x48, 0xC8, 0xFF)
GOLD   = RGBColor(0xF5, 0x9E, 0x0B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
ML = 0.75                       # left margin


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def tbox(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    """paras: list of {'runs':[(text,size,color,bold)], 'align':, 'space_after':, 'line':}"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(para.get("space_after", 4))
        p.space_before = Pt(para.get("space_before", 0))
        p.line_spacing = para.get("line", 1.1)
        for (text, size, color, bold) in para["runs"]:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Arial"
    return tb


def box(s, x, y, w, h, fill=CARD, line=BORDER, line_w=1.0, radius=0.06):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def kicker(s, text):
    tbox(s, ML, 0.55, 11.5, 0.4, [{"runs": [(text.upper(), 11.5, ACC, True)]}])


def title(s, parts, y=1.0, size=34):
    """parts: list of (text, color, bold)"""
    tbox(s, ML, y, 11.8, 1.5, [{"runs": [(t, size, c, b) for (t, c, b) in parts], "line": 1.04}])


def card(s, x, y, w, h, heading, body_lines, head_color=WHITE, fill=CARD, line=BORDER):
    box(s, x, y, w, h, fill=fill, line=line)
    paras = [{"runs": [(heading, 14.5, head_color, True)], "space_after": 7}]
    for ln in body_lines:
        paras.append({"runs": [("•  ", 11, DIM, False), (ln, 11.5, SOFT, False)], "space_after": 3, "line": 1.15})
    tbox(s, x + 0.28, y + 0.22, w - 0.55, h - 0.4, paras)


def hbar(s, x, y, label, frac, value, color, lw=2.1, bw=4.4, bh=0.30, vcolor=None):
    tbox(s, x, y - 0.05, lw, 0.42, [{"runs": [(label, 12, SOFT, True)], "align": PP_ALIGN.RIGHT}],
         anchor=MSO_ANCHOR.MIDDLE)
    bx = x + lw + 0.2
    box(s, bx, y, bw, bh, fill=TRACK, line=BORDER, radius=0.5)
    box(s, bx, y, max(0.10, bw * frac), bh, fill=color, line=None, radius=0.5)
    tbox(s, bx + bw + 0.12, y - 0.05, 1.0, 0.42, [{"runs": [(value, 12.5, vcolor or color, True)]}],
         anchor=MSO_ANCHOR.MIDDLE)


def foot(s, n):
    tbox(s, ML, 7.02, 9, 0.35, [{"runs": [("SaaSBench  ·  github.com/Bennyoooo/firmbench", 10, DIM, False)]}])
    tbox(s, 11.6, 7.02, 1.1, 0.35, [{"runs": [(f"{n:02d} / 08", 10, DIM, False)], "align": PP_ALIGN.RIGHT}])


# ============================ 1 · TITLE ============================
s = slide()
kicker(s, "HUD Frontier / RSI · RL Environments Hackathon — HUD W25 × YC")
tbox(s, ML, 1.25, 11.8, 1.4, [{"runs": [("SaaSBench", 68, WHITE, True)]}])
tbox(s, ML, 2.55, 8.6, 1.4, [{"runs": [
    ("A verifiable RL environment where an agent — or a team of agents — runs a SaaS company against ", 19, SOFT, False),
    ("up to 500,000 heterogeneous customers, graded entirely on profit.", 19, GRN, True)], "line": 1.25}])
stats = [("+260%", GRN, "held-out profit-efficiency after GRPO on an open 8B model (0.147 → 0.529)"),
         ("7%", WHITE, "of an achievable expert is all the best scripted heuristic reaches — wide-open RL headroom"),
         ("0", WHITE, "LLM judges — profit is computed exactly by a 500,000-user market sim")]
bw = 3.78
for i, (v, c, k) in enumerate(stats):
    x = ML + i * (bw + 0.18)
    box(s, x, 4.15, bw, 1.55)
    tbox(s, x + 0.25, 4.32, bw - 0.5, 0.7, [{"runs": [(v, 30, c, True)]}])
    tbox(s, x + 0.25, 5.0, bw - 0.5, 0.65, [{"runs": [(k, 10.5, DIM, False)], "line": 1.12}])
tbox(s, ML, 5.95, 11.8, 0.5, [{"runs": [
    ("Verifiable & ungameable     Domain-randomized     Trainable — the curve bends     Live on HUD", 12, SOFT, True)]}])
foot(s, 1)

# ============================ 2 · THE GAP ============================
s = slide()
kicker(s, "01 · The gap")
title(s, [("RL has math and code. It doesn't have ", WHITE, True), ("business.", GRN, True)])
tbox(s, ML, 2.0, 11.6, 1.2, [{"runs": [
    ("Frontier RL rewards provable answers — theorems, unit tests, won games. The most valuable agentic "
     "skill of all — running a business — has no verifiable gym. Grading “good strategy” with an LLM "
     "judge is noisy and gameable.", 16, SOFT, False)], "line": 1.4}])
box(s, ML, 3.35, 11.8, 1.15, fill=CARD, line=ACC)
tbox(s, ML + 0.3, 3.5, 11.2, 0.9, [{"runs": [
    ("💡  Profit is the perfect reward — a single number a simulator computes ", 17, TEXT, True),
    ("exactly", 17, GRN, True),
    (". Impossible to fake, and it only goes up if the agent truly understands the market.", 17, TEXT, True)], "line": 1.3}])
cards = [("🎯 Economically real", ["The task is the job — discovery, allocation, pricing under uncertainty."]),
         ("🔒 Verifiable", ["Reward computed by a deterministic market sim. No judge to persuade."]),
         ("🧭 Open-ended", ["Hidden latent structure forces genuine explore-then-exploit."])]
cw = 3.78
for i, (h, b) in enumerate(cards):
    card(s, ML + i * (cw + 0.18), 4.75, cw, 1.5, h, b)
foot(s, 2)

# ============================ 3 · REALISM ============================
s = slide()
kicker(s, "02 · The environment")
title(s, [("Not a demand curve — a market that ", WHITE, True), ("behaves like a real one.", GRN, True)], size=30)
tbox(s, ML, 1.95, 11.6, 0.6, [{"runs": [
    ("An agent — or a team — operates a SaaS against up to 500,000 heterogeneous customers, hidden inside a "
     "fresh randomized market each episode, learned over ~16 rounds.", 14.5, SOFT, False)], "line": 1.3}])
card(s, ML, 2.95, 5.8, 3.5, "👥  Users — a population, not a curve",
     ["their pains", "willingness-to-pay", "price elasticity", "reachable channel",
      "quality bar", "loyalty / churn", "+ per-user noise — no two alike"])
card(s, ML + 6.0, 2.95, 5.8, 3.5, "🛠️  Actions — an operator's toolkit, not a price dial",
     ["run user research", "choose pains to target", "write ad copy", "allocate budget",
      "set pricing", "build features via product specs", "manage retention"])
foot(s, 3)

# ============================ 4 · VERIFIABLE ============================
s = slide()
kicker(s, "03 · Why it's a real benchmark")
title(s, [("Execution-based reward. ", WHITE, True), ("Nothing to fake.", GRN, True)])
v = [("⚖️ No LLM judge", ["The sim computes profit from a 500k-user funnel. Can't argue your way to a score."]),
     ("📐 Honest [0,1]", ["reward = profit ÷ a theoretical ceiling — a true upper bound, not a policy."]),
     ("🎲 Domain-randomized", ["New hidden world each episode. Generalization on held-out seeds."]),
     ("🧩 Matched-quad", ["Every latent has an action, an observation, and a reward — learnable + graded."])]
cw = 2.83
for i, (h, b) in enumerate(v):
    card(s, ML + i * (cw + 0.12), 2.05, cw, 1.95, h, b)
box(s, ML, 4.25, 11.8, 1.95)
tbox(s, ML + 0.35, 4.5, 2.4, 1.4, [{"runs": [("6 / 7", 40, GRN, True)]},
                                   {"runs": [("ablation configs PASS", 12, DIM, False)]}], anchor=MSO_ANCHOR.MIDDLE)
tbox(s, ML + 3.0, 4.5, 8.4, 1.45, [
    {"runs": [("The ablation gate proves every latent is discoverable", 15, WHITE, True)], "space_after": 6},
    {"runs": [("A scripted experimenter must beat a naïve baseline and stay under an informed oracle "
               "(naive < scripted < oracle) — for v1, every single latent, and the full stack. If a latent "
               "ever became unobservable, the gate fails loudly, by construction.", 12.5, SOFT, False)], "line": 1.3}],
     anchor=MSO_ANCHOR.MIDDLE)
foot(s, 4)

# ============================ 5 · HEADROOM ============================
s = slide()
kicker(s, "04 · Headroom — it's hard")
title(s, [("A scripted expert captures ", WHITE, True), ("7% ", GRN, True), ("of what's achievable.", WHITE, True)], size=30)
tbox(s, ML, 1.85, 11.6, 0.5, [{"runs": [
    ("10 held-out seeds, full market · reward = profit ÷ theoretical ceiling.", 13.5, SOFT, False)]}])
# left: policy bars
for i, (lab, frac, val, col) in enumerate([
        ("naive", 0.012, "1.2%", DIM), ("scripted expert", 0.037, "3.7%", ACC),
        ("oracle (achievable)", 0.557, "55.7%", BLUE), ("theoretical ceiling", 1.0, "100%", GRN)]):
    hbar(s, ML, 2.55 + i * 0.52, lab, frac, val, col, lw=1.9, bw=3.2, vcolor=col)
# right: frontier models
box(s, 7.4, 2.4, 5.2, 2.55)
tbox(s, 7.65, 2.52, 4.8, 0.5, [{"runs": [("It also separates frontier models", 13.5, WHITE, True)]}])
tbox(s, 7.65, 2.9, 4.8, 0.3, [{"runs": [("seed 42 · % of ceiling", 10.5, DIM, False)]}])
for i, (lab, frac, val, col) in enumerate([
        ("GPT-5.5", 0.477, "47.7%", ACC), ("GPT-5", 0.289, "28.9%", ACC), ("Gemini 3.5", 0.112, "11.2%", ACC),
        ("Qwen3-8B · RL", 0.05, "5.0%", GRN), ("Claude Sonnet", 0.047, "4.7%", ACC), ("Opus 4.8", 0.018, "1.8%", ACC)]):
    hbar(s, 7.55, 3.25 + i * 0.27, lab, frac, val, col, lw=1.4, bw=1.9, bh=0.18, vcolor=col)
# key finding callout
box(s, ML, 5.35, 11.8, 1.05, fill=CARD, line=ACC)
tbox(s, ML + 0.3, 5.5, 11.2, 0.8, [{"runs": [
    ("Key finding:  ", 14, GRN, True),
    ("product-market fit dominates — agents that under-invest in learning the market fail, even with strong "
     "pricing or marketing. Discovery is the binding skill.", 14, TEXT, False)], "line": 1.3}])
foot(s, 5)

# ============================ 6 · THE CURVE BENDS ============================
s = slide()
kicker(s, "05 · The proof — real RL, not SFT")
title(s, [("We trained an 8B model to run a better firm — and it ", WHITE, True), ("generalized.", GRN, True)], size=28)
# left card: training curve (as points)
box(s, ML, 2.1, 5.8, 3.5)
tbox(s, ML + 0.3, 2.28, 5.2, 0.4, [{"runs": [("GRPO training — reward / epoch (in-distribution)", 12.5, SOFT, True)]}])
for i, (lab, frac, val, col) in enumerate([("epoch 0", 0.32, "0.193", ACC),
                                           ("epoch 1", 0.51, "0.307", BLUE),
                                           ("epoch 2", 0.61, "0.367", GRN)]):
    hbar(s, ML + 0.1, 2.95 + i * 0.55, lab, frac, val, col, lw=1.3, bw=3.3, vcolor=col)
tbox(s, ML + 0.3, 4.95, 5.2, 0.5, [{"runs": [("Monotonic  0.193 → 0.367  ", 13, TEXT, True), ("(+91%)", 13, GRN, True)]}])
# right card: held-out generalization
box(s, ML + 6.0, 2.1, 5.8, 3.5)
tbox(s, ML + 6.3, 2.28, 5.2, 0.4, [{"runs": [("Held-out generalization — 16 unseen worlds", 12.5, SOFT, True)]}])
hbar(s, ML + 6.1, 3.0, "base 8B", 0.245, "0.147", ACC, lw=1.5, bw=2.8, bh=0.34, vcolor=SOFT)
hbar(s, ML + 6.1, 3.6, "after GRPO", 0.88, "0.529", GRN, lw=1.5, bw=2.8, bh=0.34, vcolor=GRN)
tbox(s, ML + 6.3, 4.25, 5.2, 1.0, [{"runs": [("+260%", 40, GRN, True)]},
                                   {"runs": [("on worlds it never trained on", 12, DIM, False)]}])
box(s, ML, 5.9, 11.8, 0.95, fill=CARD, line=BORDER)
tbox(s, ML + 0.3, 6.04, 11.2, 0.7, [{"runs": [
    ("Why it's RL, not SFT:  ", 13, GRN, True),
    ("GRPO on Fireworks — the profit verifier weights the gradient (8 candidates/prompt, KL-regularized), "
     "it doesn't just filter a dataset. Ungameable end to end.", 13, TEXT, False)], "line": 1.25}])
foot(s, 6)

# ============================ 7 · TWO MODES (MULTI-AGENT) ============================
s = slide()
kicker(s, "06 · Two modes — solo founder or full company")
title(s, [("Solo founder, or a ", WHITE, True), ("whole company.", GRN, True)])
tbox(s, ML, 1.95, 11.6, 0.7, [{"runs": [
    ("Run SaaSBench with one agent, or as a team of role-agents that each see only their slice and must "
     "communicate to win — the same execution-based profit reward drives both modes.", 14.5, SOFT, False)], "line": 1.3}])
roles = [("🛠️ Builder", ["Picks & specs features.", "Must tell the team what it built."]),
         ("📣 Marketer", ["Targets pains × channels,", "writes ad copy, spends budget."]),
         ("🏷️ Pricer", ["Sets price for conversion", "vs churn — the LTV lever."]),
         ("🧭 Coordinator", ["Allocates budget, sets", "direction, commits the round."])]
cw = 2.83
for i, (h, b) in enumerate(roles):
    card(s, ML + i * (cw + 0.12), 2.95, cw, 1.55, h, b)
card(s, ML, 4.7, 5.8, 1.6, "📉 A coordination tax — new skill, same rigor",
     ["Roles see partial views & must share what they learn.",
      "Grade team profit vs a single-agent full-info oracle — the gap is the cost of poor coordination.",
      "Still execution-based, still ungameable."], fill=CARD, line=ACC)
card(s, ML + 6.0, 4.7, 5.8, 1.6, "🧠 One checkpoint, four roles",
     ["Not four models — one shared policy, role-conditioned by prompt.",
      "Every role-turn shares the team's profit reward; GRPO over the pooled turns.",
      "Parameter sharing — the only sane shape for cooperative LLM agents."])
foot(s, 7)

# ============================ 8 · SHIPPED + CLOSE ============================
s = slide()
kicker(s, "07 · Shipped — runnable today")
title(s, [("Deployed on HUD. ", WHITE, True), ("Run it now.", GRN, True)])
card(s, ML, 2.05, 5.8, 1.9, "🚀 Live environment",
     ["Full env on HUD with 6 MCP tools.", "Evaluate any frontier model — or our fine-tuned checkpoint.",
      "hud eval tasks.py claude --task-ids market_discovery_seed42"])
card(s, ML + 6.0, 2.05, 5.8, 1.9, "🎬 Inspect every run",
     ["Replay viewer steps through any episode — campaigns, features, profit curve.",
      "Model leaderboard + a world explorer for the hidden market."])
tbox(s, ML, 4.2, 11.6, 0.4, [{"runs": [("BUILT ON THE SPONSOR STACK", 11, DIM, True)]}])
tbox(s, ML, 4.55, 11.6, 0.5, [{"runs": [
    ("HUD — host + on-policy RL (hud.train)     Fireworks — inference + GRPO     OpenAI · Anthropic · Google — leaderboard",
     12.5, SOFT, True)], "line": 1.2}])
box(s, ML, 5.25, 11.8, 1.45, fill=CARD, line=ACC)
tbox(s, ML, 5.45, 11.8, 1.05, [
    {"runs": [("A benchmark where “make the number go up” means “build a better business.”", 20, WHITE, True)],
     "align": PP_ALIGN.CENTER, "space_after": 6},
    {"runs": [("Verifiable. Trainable. And far from solved.", 14, GRN, True)], "align": PP_ALIGN.CENTER}],
     anchor=MSO_ANCHOR.MIDDLE)
foot(s, 8)

prs.save("SaaSBench.pptx")
print("wrote SaaSBench.pptx with", len(prs.slides._sldIdLst), "slides")
